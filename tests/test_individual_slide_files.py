"""
Integration tests for individual slide file storage feature.

Tests that slides are stored as individual .pptx files during ingestion
and can be used for assembly without dependency on original file paths.
"""

import pytest
import tempfile
from pathlib import Path
from pptx import Presentation

from slidex.core.ingest import ingest_engine
from slidex.core.assembler import slide_assembler
from slidex.core.database import db
from slidex.config import settings


@pytest.fixture
def temp_presentation():
    """Create a temporary test presentation with 2 slides."""
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        prs = Presentation()
        
        # Add first slide
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.shapes.title.text = 'Integration Test Slide 1'
        slide1.shapes.placeholders[1].text = 'Content for integration test slide 1'
        
        # Add second slide
        slide2 = prs.slides.add_slide(prs.slide_layouts[0])
        slide2.shapes.title.text = 'Integration Test Slide 2'
        slide2.shapes.placeholders[1].text = 'Content for integration test slide 2'
        
        prs.save(tmp.name)
        tmp_path = Path(tmp.name)
    
    yield tmp_path
    
    # Cleanup
    if tmp_path.exists():
        tmp_path.unlink()


def test_individual_slide_files_created(temp_presentation):
    """Test that individual slide files are created during ingestion."""
    # Ingest the test presentation
    deck_id = ingest_engine.ingest_file(
        temp_presentation,
        uploader="integration_test",
        session_id="test_session"
    )
    
    assert deck_id is not None, "Ingestion should return a deck_id"
    
    # Get slides from database
    slides = db.get_slides_by_deck_id(deck_id)
    assert len(slides) == 2, "Should have 2 slides"
    
    # Verify each slide has a slide_file_path
    for slide in slides:
        assert slide['slide_file_path'] is not None, "Each slide should have slide_file_path"
        
        # Verify the file exists
        slide_file = Path(slide['slide_file_path'])
        assert slide_file.exists(), f"Slide file should exist: {slide_file}"
        
        # Verify it's a valid PowerPoint file with exactly 1 slide
        slide_prs = Presentation(str(slide_file))
        assert len(slide_prs.slides) == 1, "Individual slide file should contain exactly 1 slide"
        
        # Verify the slide has content
        assert slide_prs.slides[0].shapes.title.text in [
            'Integration Test Slide 1',
            'Integration Test Slide 2'
        ], "Slide should have correct title"


def test_assembly_uses_individual_files(temp_presentation):
    """Test that assembly can use individual slide files."""
    # Ingest the test presentation
    deck_id = ingest_engine.ingest_file(
        temp_presentation,
        uploader="integration_test",
        session_id="test_session"
    )
    
    # Get slide IDs
    slides = db.get_slides_by_deck_id(deck_id)
    slide_ids = [slide['slide_id'] for slide in slides]
    
    # Assemble presentation
    output_path = slide_assembler.assemble(
        slide_ids=slide_ids,
        output_filename="test_integration_assembled.pptx",
        preserve_order=True
    )
    
    assert output_path.exists(), "Assembled presentation should exist"
    
    # Verify assembled presentation
    assembled_prs = Presentation(str(output_path))
    assert len(assembled_prs.slides) == 2, "Assembled presentation should have 2 slides"
    
    # Verify slide content is preserved
    assert assembled_prs.slides[0].shapes.title.text == 'Integration Test Slide 1'
    assert assembled_prs.slides[1].shapes.title.text == 'Integration Test Slide 2'
    
    # Cleanup
    if output_path.exists():
        output_path.unlink()


def test_individual_slides_are_portable(temp_presentation):
    """Test that individual slide files work independently of original file."""
    # Ingest the test presentation
    deck_id = ingest_engine.ingest_file(
        temp_presentation,
        uploader="integration_test",
        session_id="test_session"
    )
    
    # Get first slide
    slides = db.get_slides_by_deck_id(deck_id)
    first_slide = slides[0]
    slide_file_path = Path(first_slide['slide_file_path'])
    
    # Delete the original presentation file
    temp_presentation.unlink()
    assert not temp_presentation.exists(), "Original file should be deleted"
    
    # Verify individual slide file still works
    assert slide_file_path.exists(), "Individual slide file should still exist"
    
    slide_prs = Presentation(str(slide_file_path))
    assert len(slide_prs.slides) == 1, "Individual slide file should work independently"
    assert slide_prs.slides[0].shapes.title.text == 'Integration Test Slide 1'


def test_slide_file_naming_convention():
    """Test that slide files are named using slide_id."""
    # Slide files should be named: {slide_id}.pptx
    # This ensures uniqueness and allows easy lookup
    
    slides_dir = settings.slides_dir
    assert slides_dir.exists(), "Slides directory should exist"
    
    # Check existing slide files follow naming convention
    slide_files = list(slides_dir.glob("*.pptx"))
    
    for slide_file in slide_files:
        # Filename (without extension) should be a valid UUID
        slide_id = slide_file.stem
        assert len(slide_id) == 36, f"Slide ID should be UUID format: {slide_id}"
        assert slide_id.count('-') == 4, f"Slide ID should have 4 hyphens: {slide_id}"

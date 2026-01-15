"""
Slide assembler for creating new PowerPoint presentations from selected slides.
"""

from pathlib import Path
from typing import List, Optional
from datetime import datetime
from pptx import Presentation
from pptx.parts.image import Image, ImagePart
from pptx.opc.package import Part
from copy import deepcopy
from pptx.oxml.ns import qn
import io

from slidex.config import settings
from slidex.logging_config import logger
from slidex.core.database import db


class SlideAssembler:
    """Assembles selected slides into a new PowerPoint presentation."""
    
    @staticmethod
    def assemble(
        slide_ids: List[str],
        output_filename: Optional[str] = None,
        preserve_order: bool = False,
    ) -> Path:
        """Assemble slides into a new PowerPoint presentation.

        Args:
            slide_ids: List of slide IDs to include
            output_filename: Optional output filename (generated if not provided)
            preserve_order: If True, preserve the order of slide_ids;
                            otherwise order by original deck order
        """
        if not slide_ids:
            raise ValueError("No slides provided for assembly")

        logger.info(f"Assembling presentation with {len(slide_ids)} slides (PPTX)")
        
        # Fetch slide metadata
        slides_data = []
        for slide_id in slide_ids:
            slide = db.get_slide_by_id(slide_id)
            if not slide:
                logger.warning(f"Slide not found: {slide_id}, skipping")
                continue
            slides_data.append(slide)
        
        if not slides_data:
            raise ValueError("No valid slides found")
        
        # Determine the most common slide dimensions to use for the output
        dimensions_count = {}
        for slide_data in slides_data:
            slide_file_path = slide_data.get('slide_file_path')
            try:
                if slide_file_path and Path(slide_file_path).exists():
                    temp_prs = Presentation(slide_file_path)
                else:
                    deck_path = slide_data['deck_path']
                    temp_prs = Presentation(deck_path)
                
                dim = (temp_prs.slide_width, temp_prs.slide_height)
                dimensions_count[dim] = dimensions_count.get(dim, 0) + 1
            except Exception:
                pass
        
        # Use the most common dimensions, or default to 16:9 widescreen
        if dimensions_count:
            target_dimensions = max(dimensions_count.items(), key=lambda x: x[1])[0]
        else:
            # Default to 16:9 widescreen (10 inches x 5.625 inches)
            target_dimensions = (9144000, 5143500)
        
        logger.info(f"Using slide dimensions: {target_dimensions[0]} x {target_dimensions[1]}")
        if len(dimensions_count) > 1:
            logger.warning(f"Mixed slide dimensions detected: {dimensions_count}. Normalizing to most common size.")
        
        # Create new presentation with target dimensions
        new_prs = Presentation()
        new_prs.slide_width = target_dimensions[0]
        new_prs.slide_height = target_dimensions[1]
        
        # Counter for unique image naming (no deduplication)
        image_counter = {'count': 0}
        
        # Process slides in the requested order
        if preserve_order:
            ordered_slides = slides_data
        else:
            # Order by deck path and then slide index
            ordered_slides = sorted(
                slides_data,
                key=lambda x: (x['deck_path'], x['slide_index'])
            )
        
        for slide_data in ordered_slides:
            slide_file_path = slide_data.get('slide_file_path')
            
            try:
                # If individual slide file exists, use it (new method)
                if slide_file_path and Path(slide_file_path).exists():
                    # Load the individual slide file (contains just one slide)
                    source_prs = Presentation(slide_file_path)
                    source_slide = source_prs.slides[0]  # Always first slide
                    logger.debug(f"Loaded individual slide file: {slide_file_path}")
                else:
                    # Fallback to old method: load from original deck
                    deck_path = slide_data['deck_path']
                    slide_index = slide_data['slide_index']
                    source_prs = Presentation(deck_path)
                    source_slide = source_prs.slides[slide_index]
                    logger.debug(f"Loaded slide from original deck: {deck_path} (index {slide_index})")
                
                # Copy slide to new presentation with unique image naming
                SlideAssembler._copy_slide(source_slide, new_prs, image_counter)
                
                logger.debug(
                    f"Copied slide: {slide_data.get('title_header') or 'Untitled'}"
                )
                
            except Exception as e:
                logger.error(
                    f"Error copying slide {slide_data.get('slide_id')}: {e}"
                )
                # Continue with other slides
                continue
        
        # Generate output filename if not provided
        if not output_filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_filename = f"assembled_{timestamp}.pptx"
        
        # Ensure output directory exists
        output_path = settings.exports_dir / output_filename
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        # Save new presentation
        new_prs.save(str(output_path))
        
        logger.info(f"Presentation assembled: {output_path} ({len(new_prs.slides)} slides)")
        
        return output_path
    
    @staticmethod
    def _copy_slide(source_slide, target_presentation: Presentation, image_counter: dict = None) -> None:
        """
        Copy a slide from source to target presentation using deep XML cloning.
        This preserves all formatting, shapes, images, and content.
        
        Args:
            source_slide: Source slide to copy
            target_presentation: Target presentation
            image_counter: Optional dict with 'count' key for unique image naming
        
        Note: Target presentation dimensions should already be set before calling this.
        """
        if image_counter is None:
            image_counter = {'count': 0}
        
        # Get a blank slide layout from target
        blank_layout = target_presentation.slide_layouts[6]

        # Add new slide with blank layout
        new_slide = target_presentation.slides.add_slide(blank_layout)

        # Deep copy the entire slide XML tree from source to target
        # Collect newly inserted elements so we can remap relationship ids
        new_elements = []
        for shape in source_slide.shapes:
            try:
                el = shape.element
                newel = deepcopy(el)
                new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
                new_elements.append(newel)
            except Exception as e:
                logger.debug(f"Could not copy shape {getattr(shape, 'name', '')}: {e}")

        # Copy all relationships including images, media, and embedded objects (OLE, charts, etc.)
        # Map old rIds -> new rIds so we can remap r:embed and r:id attributes in the copied XML
        old_to_new = {}
        try:
            for rId, rel in list(source_slide.part.rels.items()):
                try:
                    related_part = rel.target_part
                    new_part = None
                    
                    # Handle different relationship types
                    reltype_lower = rel.reltype.lower()
                    
                    # Skip internal PPTX structure relationships (layouts, masters, themes)
                    # These should not be copied - they belong to the presentation template
                    skip_reltypes = ['slidelayout', 'slidemaster', 'theme', 'notesmaster', 'notesslide', 'handoutmaster']
                    if any(skip in reltype_lower for skip in skip_reltypes):
                        logger.debug(f"Skipping internal relationship: {rel.reltype}")
                        continue
                    
                    if 'image' in reltype_lower:
                        # Handle images
                        try:
                            blob = getattr(related_part, 'blob', None)
                            content_type = getattr(related_part, 'content_type', None)
                            target_pkg = target_presentation.part.package

                            if blob is not None and content_type:
                                try:
                                    image_counter['count'] += 1
                                    image = Image.from_blob(blob)
                                    new_part = ImagePart.new(target_pkg, image)
                                except Exception as e:
                                    logger.debug(f"Could not create image part: {e}")
                                    new_part = None
                        except Exception as e:
                            logger.debug(f"Could not duplicate image: {e}")
                    
                    elif 'oleobject' in reltype_lower or 'package' in reltype_lower or 'embeddings' in reltype_lower:
                        # Handle OLE objects, embedded packages, and smart objects
                        # These need special handling - copy the blob as-is
                        try:
                            blob = getattr(related_part, 'blob', None)
                            if blob is not None:
                                # Create a new part in the target package with the same content type
                                content_type = getattr(related_part, 'content_type', 'application/vnd.openxmlformats-officedocument.oleObject')
                                
                                # Generate unique part name
                                image_counter['count'] += 1
                                ext = 'bin'  # default extension
                                if hasattr(related_part, 'partname'):
                                    original_ext = str(related_part.partname).split('.')[-1]
                                    if original_ext:
                                        ext = original_ext
                                
                                unique_partname = f'/ppt/embeddings/oleObject{image_counter["count"]}.{ext}'
                                
                                # Add part directly to package
                                new_part = Part(unique_partname, content_type, blob, target_presentation.part.package)
                                target_presentation.part.package._parts[unique_partname] = new_part
                                
                                logger.debug(f"Copied OLE/embedded object: {unique_partname}")
                        except Exception as e:
                            logger.debug(f"Could not copy OLE object: {e}")
                    
                    elif 'media' in reltype_lower or 'video' in reltype_lower or 'audio' in reltype_lower:
                        # Handle media files (video, audio)
                        try:
                            blob = getattr(related_part, 'blob', None)
                            if blob is not None:
                                content_type = getattr(related_part, 'content_type', 'application/octet-stream')
                                image_counter['count'] += 1
                                
                                # Determine extension
                                ext = 'bin'
                                if hasattr(related_part, 'partname'):
                                    original_ext = str(related_part.partname).split('.')[-1]
                                    if original_ext:
                                        ext = original_ext
                                
                                unique_partname = f'/ppt/media/media{image_counter["count"]}.{ext}'
                                
                                new_part = Part(unique_partname, content_type, blob, target_presentation.part.package)
                                target_presentation.part.package._parts[unique_partname] = new_part
                                
                                logger.debug(f"Copied media: {unique_partname}")
                        except Exception as e:
                            logger.debug(f"Could not copy media: {e}")
                    
                    else:
                        # For other relationship types (charts, hyperlinks, etc.), try generic copy
                        try:
                            blob = getattr(related_part, 'blob', None)
                            if blob is not None:
                                content_type = getattr(related_part, 'content_type', 'application/octet-stream')
                                image_counter['count'] += 1
                                
                                ext = 'bin'
                                if hasattr(related_part, 'partname'):
                                    original_ext = str(related_part.partname).split('.')[-1]
                                    if original_ext:
                                        ext = original_ext
                                
                                unique_partname = f'/ppt/other/object{image_counter["count"]}.{ext}'
                                
                                new_part = Part(unique_partname, content_type, blob, target_presentation.part.package)
                                target_presentation.part.package._parts[unique_partname] = new_part
                                
                                logger.debug(f"Copied other relationship type {rel.reltype}: {unique_partname}")
                        except Exception as e:
                            logger.debug(f"Could not copy relationship {rel.reltype}: {e}")

                    # Create relationship to the new part (or fallback to original)
                    try:
                        if new_part is not None:
                            new_rel = new_slide.part.relate_to(new_part, rel.reltype)
                        else:
                            # Last-resort: relate to original part
                            new_rel = new_slide.part.relate_to(related_part, rel.reltype)

                        old_to_new[rId] = getattr(new_rel, 'rId', None)
                    except Exception as e:
                        logger.debug(f"Could not create relationship on new slide: {e}")
                        
                except Exception as e:
                    logger.debug(f"Skipping relationship {rId}: {e}")
        except Exception as e:
            logger.debug(f"Error copying slide relationships: {e}")

        # Remap r:embed and r:id references inside the newly inserted XML elements
        if old_to_new:
            try:
                for newel in new_elements:
                    for node in newel.iter():
                        # Remap r:embed attributes (e.g. in a:blip for images)
                        embed_val = node.get(qn('r:embed'))
                        if embed_val and embed_val in old_to_new:
                            new_rid = old_to_new.get(embed_val)
                            if new_rid:
                                node.set(qn('r:embed'), new_rid)
                        
                        # Remap r:id attributes (e.g. in oleObject, media references)
                        id_val = node.get(qn('r:id'))
                        if id_val and id_val in old_to_new:
                            new_rid = old_to_new.get(id_val)
                            if new_rid:
                                node.set(qn('r:id'), new_rid)
            except Exception as e:
                logger.debug(f"Error remapping relationship attributes: {e}")
    


# Global assembler instance
slide_assembler = SlideAssembler()

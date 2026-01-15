"""
Slide processor for extracting text and generating thumbnails from PowerPoint slides.
"""

import io
import subprocess
import tempfile
from copy import deepcopy
from pathlib import Path
from typing import Optional, Tuple
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

from slidex.config import settings
from slidex.logging_config import logger


class SlideProcessor:
    """Process PowerPoint slides to extract text and generate thumbnails."""
    
    @staticmethod
    def extract_text_from_shape(shape) -> str:
        """Extract text from a single shape."""
        if hasattr(shape, "text"):
            return shape.text
        return ""
    
    @staticmethod
    def extract_visual_content_info(slide) -> str:
        """
        Extract information about visual content in a slide.
        
        Returns:
            Description of visual elements (images, charts, tables, shapes)
        """
        visual_elements = []
        
        if hasattr(slide, 'shapes'):
            image_count = 0
            chart_count = 0
            table_count = 0
            shape_count = 0
            
            for shape in slide.shapes:
                try:
                    # Count pictures
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_count += 1
                    # Count charts
                    elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                        chart_count += 1
                    # Count tables - check shape_type first to avoid hasattr triggering ValueError
                    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        table_count += 1
                    # Count other shapes (rectangles, circles, etc.)
                    elif shape.shape_type in [
                        MSO_SHAPE_TYPE.AUTO_SHAPE,
                        MSO_SHAPE_TYPE.FREEFORM,
                        MSO_SHAPE_TYPE.GROUP
                    ]:
                        shape_count += 1
                except AttributeError:
                    continue
            
            # Build description
            if image_count > 0:
                visual_elements.append(f"{image_count} image(s)")
            if chart_count > 0:
                visual_elements.append(f"{chart_count} chart(s)")
            if table_count > 0:
                visual_elements.append(f"{table_count} table(s)")
            if shape_count > 0:
                visual_elements.append(f"{shape_count} shape(s)")
        
        if visual_elements:
            return "Contains: " + ", ".join(visual_elements)
        return ""
    
    @staticmethod
    def extract_text_from_slide(slide) -> Tuple[Optional[str], str]:
        """
        Extract text from a slide.
        
        Returns:
            Tuple of (title_header, plain_text)
        """
        title_header = None
        text_parts = []
        
        # Try to get title from slide
        if hasattr(slide, 'shapes'):
            for shape in slide.shapes:
                # Check if this is a title shape
                if hasattr(shape, 'shape_type') and shape.name.lower().startswith('title'):
                    if hasattr(shape, 'text') and shape.text.strip():
                        title_header = shape.text.strip()
                
                # Extract all text
                text = SlideProcessor.extract_text_from_shape(shape)
                if text.strip():
                    text_parts.append(text.strip())
                
                # Extract text from tables
                # Note: shape.table raises ValueError if shape doesn't contain a table
                    try:
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    text_parts.append(cell.text.strip())
                    except (ValueError, AttributeError):
                        # Shape doesn't contain a table, skip
                        pass
        
        # Extract notes
        if hasattr(slide, 'notes_slide') and slide.notes_slide:
            if hasattr(slide.notes_slide, 'notes_text_frame'):
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text.strip():
                    text_parts.append(f"[Notes: {notes_text.strip()}]")
        
        plain_text = "\n".join(text_parts)
        
        # If no title found, try to use first non-empty text as title
        if not title_header and text_parts:
            title_header = text_parts[0][:100]  # Use first 100 chars
        
        return title_header, plain_text
    
    @staticmethod
    def generate_thumbnail(
        presentation: Presentation,
        slide_index: int,
        output_path: Path,
        width: int = 320,
        deck_pdf_path: Optional[Path] = None,
    ) -> None:
        """
        Generate a thumbnail image for a slide.

        Preferred strategy:
        1. If a deck-level PDF is available, render the corresponding PDF page
           for a pixel-perfect preview.
        2. Fall back to enhanced Pillow rendering if PDF rendering is not
           available or fails.
        
        Args:
            presentation: The PowerPoint presentation
            slide_index: Index of the slide
            output_path: Path to save the thumbnail
            width: Width of the thumbnail in pixels
        """
        output_path.parent.mkdir(parents=True, exist_ok=True)

        # Prefer PDF-based rendering when a deck PDF is available
        if deck_pdf_path is not None:
            try:
                from slidex.core.pdf_processor import pdf_processor

                img = pdf_processor.render_page_to_image(deck_pdf_path, slide_index, width)
                if img is not None:
                    img.save(output_path, "PNG")
                    logger.debug(
                        f"PDF-based thumbnail generated: {output_path} (page {slide_index})"
                    )
                    return
            except Exception as e:  # pragma: no cover - defensive
                logger.error(f"PDF-based thumbnail generation failed: {e}")

        # Fall back to enhanced Pillow rendering
        SlideProcessor._generate_pillow_thumbnail(presentation, slide_index, output_path, width)
    
    
    @staticmethod
    def _generate_pillow_thumbnail(
        presentation: Presentation,
        slide_index: int,
        output_path: Path,
        width: int = 320
    ) -> None:
        """
        Generate thumbnail using enhanced Pillow rendering.
        
        Renders slide backgrounds, shapes, images, and text content.
        """
        try:
            slide = presentation.slides[slide_index]
            
            # Calculate dimensions (assuming 16:9 aspect ratio)
            height = int(width * 9 / 16)
            
            # Create image with slide background
            img = Image.new('RGB', (width, height), color='white')
            
            # Try to get and render slide background
            if slide.background.fill.type:
                try:
                    bg_color = slide.background.fill.fore_color.rgb
                    bg_rgb = (bg_color[0], bg_color[1], bg_color[2])
                    img = Image.new('RGB', (width, height), color=bg_rgb)
                except:
                    pass  # Use default white background
            
            draw = ImageDraw.Draw(img)
            
            # Load fonts
            try:
                title_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 20)
                text_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 12)
            except:
                title_font = ImageFont.load_default()
                text_font = ImageFont.load_default()
            
            # Render shapes with basic styling
            y_offset = 10
            drawn_title = False
            images_drawn = 0
            
            if hasattr(slide, 'shapes'):
                for shape in slide.shapes:
                    try:
                        # Try to render images
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            try:
                                image = shape.image
                                image_bytes = image.blob
                                from PIL import Image as PILImage
                                pil_img = PILImage.open(io.BytesIO(image_bytes))
                                
                                # Scale image to fit thumbnail
                                max_img_width = width - 20
                                max_img_height = height - y_offset - 10
                                
                                pil_img.thumbnail((max_img_width, max_img_height), PILImage.Resampling.LANCZOS)
                                
                                # Paste image onto thumbnail
                                if pil_img.mode == 'RGBA':
                                    img.paste(pil_img, (10, y_offset), pil_img)
                                else:
                                    img.paste(pil_img, (10, y_offset))
                                
                                y_offset += pil_img.height + 5
                                images_drawn += 1
                                continue
                            except Exception as e:
                                logger.debug(f"Could not render image: {e}")
                        
                        # Skip text-less shapes
                        if not hasattr(shape, 'text') or not shape.text.strip():
                            continue
                        
                        # Get text color if available
                        text_color = (0, 0, 0)  # Default black
                        if hasattr(shape, 'text_frame'):
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.font.color.type:
                                        try:
                                            rgb = run.font.color.rgb
                                            text_color = (rgb[0], rgb[1], rgb[2])
                                        except:
                                            pass
                        
                        # Draw title differently
                        if hasattr(shape, 'name') and shape.name.lower().startswith('title') and not drawn_title:
                            title_wrapped = SlideProcessor._wrap_text(
                                shape.text.strip(),
                                width - 20,
                                draw,
                                title_font
                            )
                            draw.text((10, y_offset), title_wrapped, fill=text_color, font=title_font)
                            y_offset += 50
                            drawn_title = True
                        elif y_offset < height - 30:
                            # Draw regular text
                            text_wrapped = SlideProcessor._wrap_text(
                                shape.text.strip()[:200],  # Limit content
                                width - 20,
                                draw,
                                text_font
                            )
                            draw.text((10, y_offset), text_wrapped, fill=(100, 100, 100), font=text_font)
                            y_offset += 40
                    
                    except Exception as e:
                        logger.debug(f"Error rendering shape: {e}")
                        continue
            
            # Save thumbnail
            img.save(output_path, 'PNG')
            logger.debug(f"Pillow thumbnail generated: {output_path} (images: {images_drawn})")
        
        except Exception as e:
            logger.error(f"Error generating Pillow thumbnail for slide {slide_index}: {e}")
            # Create a placeholder image on error
            img = Image.new('RGB', (width, int(width * 9 / 16)), color='lightgray')
            draw = ImageDraw.Draw(img)
            try:
                draw.text((width//2 - 50, int(width * 9 / 16)//2 - 10), "No Preview", fill='black')
            except:
                pass
            output_path.parent.mkdir(parents=True, exist_ok=True)
            img.save(output_path, 'PNG')
    
    @staticmethod
    def _wrap_text(text: str, max_width: int, draw, font) -> str:
        """Simple text wrapping helper."""
        words = text.split()
        lines = []
        current_line = []
        
        for word in words:
            test_line = ' '.join(current_line + [word])
            # For default font, estimate width
            try:
                bbox = draw.textbbox((0, 0), test_line, font=font)
                width = bbox[2] - bbox[0]
            except:
                width = len(test_line) * 7  # Rough estimate
            
            if width <= max_width:
                current_line.append(word)
            else:
                if current_line:
                    lines.append(' '.join(current_line))
                current_line = [word]
        
        if current_line:
            lines.append(' '.join(current_line))
        
        return '\n'.join(lines[:10])  # Limit to 10 lines
    
    @staticmethod
    def save_slide_as_file(
        presentation: Presentation,
        slide_index: int,
        output_path: Path
    ) -> None:
        """
        Save a single slide as a standalone .pptx file.
        
        Creates a new presentation containing only the specified slide,
        preserving all formatting, shapes, images, and content.
        
        Args:
            presentation: The source PowerPoint presentation
            slide_index: Index of the slide to extract
            output_path: Path where the single-slide .pptx should be saved
        """
        try:
            output_path.parent.mkdir(parents=True, exist_ok=True)
            
            # Get the source slide
            source_slide = presentation.slides[slide_index]
            
            # Create a new presentation with just this slide
            new_prs = Presentation()
            
            # Copy slide dimensions from source
            try:
                new_prs.slide_width = presentation.slide_width
                new_prs.slide_height = presentation.slide_height
            except Exception as e:
                logger.debug(f"Could not copy slide dimensions: {e}")
            
            # Add blank slide with blank layout
            blank_layout = new_prs.slide_layouts[6]  # Blank layout
            new_slide = new_prs.slides.add_slide(blank_layout)
            
            # Build a mapping of old rIds to new rIds for relationship updates
            rel_id_map = {}
            
            # Copy image and media parts first, before copying shapes
            try:
                for rel_id, rel in source_slide.part.rels.items():
                    if 'image' in rel.reltype or 'media' in rel.reltype:
                        try:
                            # Get the image/media part and its binary content
                            related_part = rel.target_part
                            
                            # Create a new image part in the destination package
                            # Wrap blob in BytesIO since get_or_add_image_part expects a file-like object
                            from io import BytesIO
                            image_stream = BytesIO(related_part.blob)
                            image_part = new_slide.part.package.get_or_add_image_part(image_stream)
                            
                            # Create relationship in new slide to the image
                            # relate_to returns the rId string directly
                            new_rId = new_slide.part.relate_to(image_part, rel.reltype)
                            
                            # Map old rId to new rId for updating shape references
                            rel_id_map[rel_id] = new_rId
                            
                        except Exception as e:
                            logger.debug(f"Could not copy image/media part: {e}")
            except Exception as e:
                logger.debug(f"Error copying image/media relationships: {e}")
            
            # Copy all shapes from source slide using deep XML cloning
            for shape in source_slide.shapes:
                try:
                    el = shape.element
                    newel = deepcopy(el)
                    
                    # Update relationship IDs in the copied element to match new relationships
                    # This handles image references in picture shapes
                    for old_rid, new_rid in rel_id_map.items():
                        # Update embed attribute (for images)
                        for blip in newel.iter():
                            if blip.tag.endswith('}blip'):
                                embed_attr = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
                                if blip.get(embed_attr) == old_rid:
                                    blip.set(embed_attr, new_rid)
                    
                    new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
                except Exception as e:
                    logger.debug(f"Could not copy shape {getattr(shape, 'name', '')}: {e}")
            
            # Save the single-slide presentation
            new_prs.save(str(output_path))
            logger.debug(f"Saved individual slide file: {output_path}")
            
        except Exception as e:
            logger.error(f"Error saving slide as file: {e}")
            raise


# Global slide processor instance
slide_processor = SlideProcessor()

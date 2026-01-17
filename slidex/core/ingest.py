"""
Ingestion engine for processing PowerPoint files.
Handles file discovery, deduplication, text extraction, embedding generation, and storage.
"""

import hashlib
import os
from pathlib import Path
from typing import Optional, List
from pptx import Presentation

from slidex.config import settings
from slidex.logging_config import logger
from slidex.core.database import db
from slidex.core.ollama_client import ollama_client
from slidex.core.slide_processor import slide_processor
from slidex.core.lightrag_client import lightrag_client
from slidex.core.pdf_processor import pdf_processor


class IngestEngine:
    """Engine for ingesting PowerPoint presentations."""
    
    @staticmethod
    def compute_file_hash(file_path: Path) -> str:
        """Compute SHA256 hash of a file."""
        sha256_hash = hashlib.sha256()
        
        with open(file_path, "rb") as f:
            # Read in chunks to handle large files
            for byte_block in iter(lambda: f.read(4096), b""):
                sha256_hash.update(byte_block)
        
        # Include file size and mtime for extra uniqueness
        stat = file_path.stat()
        hash_input = f"{sha256_hash.hexdigest()}_{stat.st_size}_{stat.st_mtime}"
        
        return hashlib.sha256(hash_input.encode()).hexdigest()
    
    @staticmethod
    def ingest_file(
        file_path: Path,
        uploader: Optional[str] = None,
        session_id: Optional[str] = None
    ) -> Optional[str]:
        """
        Ingest a single PowerPoint file.
        
        Args:
            file_path: Path to the .pptx file
            uploader: Optional uploader name
            session_id: Optional session ID for audit logging
            
        Returns:
            deck_id if successful, None if skipped (duplicate)
        """
        file_path = Path(file_path).resolve()
        
        if not file_path.exists():
            logger.error(f"File not found: {file_path}")
            raise FileNotFoundError(f"File not found: {file_path}")
        
        if not file_path.suffix.lower() == '.pptx':
            logger.error(f"Not a PowerPoint file: {file_path}")
            raise ValueError(f"Not a PowerPoint file: {file_path}")
        
        logger.info(f"Starting ingestion: {file_path}")
        
        # Compute file hash
        file_hash = IngestEngine.compute_file_hash(file_path)
        logger.debug(f"File hash: {file_hash}")
        
        # Check if already ingested
        existing_deck_id = db.check_deck_exists(file_hash)
        if existing_deck_id:
            logger.info(f"File already ingested (deck_id: {existing_deck_id}), skipping")
            return None
        
        # Load presentation
        try:
            presentation = Presentation(str(file_path))
        except Exception as e:
            logger.error(f"Error loading presentation: {e}")
            raise
        
        slide_count = len(presentation.slides)
        logger.info(f"Loaded presentation: {slide_count} slides")

        # Convert whole deck to PDF once (if enabled) so we can reuse it for
        # thumbnails and per-slide PDFs.
        deck_pdf_path = None
        if settings.pdf_conversion_enabled:
            try:
                deck_pdf_path = pdf_processor.convert_pptx_to_pdf(file_path)
            except Exception as e:  # pragma: no cover - defensive
                logger.error(f"Deck-level PDF conversion failed: {e}")
                deck_pdf_path = None
        
        # Insert deck record
        deck_id = db.insert_deck(
            file_hash=file_hash,
            original_path=str(file_path),
            filename=file_path.name,
            slide_count=slide_count,
            uploader=uploader,
        )
        
        # Collect slide data for batch insertion into LightRAG
        lightrag_documents = []

        # Collect slide metadata for PDF extraction and DB insert
        slide_data_for_pdf = []
        
        # Process each slide
        for slide_idx, slide in enumerate(presentation.slides):
            logger.debug(f"Processing slide {slide_idx + 1}/{slide_count}")
            
            # Extract text
            title_header, plain_text = slide_processor.extract_text_from_slide(slide)
            
            # Extract visual content information
            visual_context = slide_processor.extract_visual_content_info(slide)
            
            if not plain_text.strip():
                logger.warning(f"Slide {slide_idx} has no text, using placeholder")
                plain_text = f"[Slide {slide_idx + 1}]"
            
            # Generate slide_id early so we can use it for file naming
            slide_id = db.generate_slide_id()
            
            # Save individual slide as standalone .pptx file
            slide_filename = f"{slide_id}.pptx"
            slide_file_path = settings.slides_dir / slide_filename
            slide_processor.save_slide_as_file(
                presentation,
                slide_idx,
                slide_file_path
            )
            
            # Generate thumbnail
            thumbnail_filename = f"{deck_id}_{slide_idx}.png"
            thumbnail_path = settings.thumbnails_dir / deck_id / thumbnail_filename
            slide_processor.generate_thumbnail(
                presentation,
                slide_idx,
                thumbnail_path,
                width=settings.thumbnail_width,
                deck_pdf_path=deck_pdf_path,
            )
            
            # Generate summary
            try:
                # For very minimal text (< 15 chars), create simple summary without LLM
                if len(plain_text.strip()) < 15:
                    if visual_context:
                        summary = f"{plain_text.strip()} - {visual_context}"
                    else:
                        summary = f"Slide: {plain_text.strip()}"
                    logger.debug(f"Using simple summary for minimal text slide")
                else:
                    summary = ollama_client.generate_summary(
                        plain_text,
                        max_words=20,
                        session_id=session_id,
                        visual_context=visual_context if visual_context else None
                    )
            except Exception as e:
                logger.error(f"Error generating summary for slide {slide_idx}: {e}")
                summary = plain_text[:100]  # Fallback to truncated text
            
            # Complexity and PDF preference are no longer used for control flow.
            # Keep default values for backward-compatible schema.
            complexity_score = 0
            requires_pdf = False

            # Prepare embedding input
            embedding_input = f"{title_header or ''}\n{plain_text}\n{summary}"
            
            # Placeholder vector_id (not used with LightRAG)
            vector_id = slide_idx
            
            # Store paths relative to project root, or absolute if outside
            try:
                # Try to make it relative to current working directory
                rel_thumbnail_path = thumbnail_path.resolve().relative_to(Path.cwd().resolve())
                thumbnail_path_str = str(rel_thumbnail_path)
            except ValueError:
                # If not in subpath, use absolute path
                thumbnail_path_str = str(thumbnail_path.resolve())
            
            try:
                rel_slide_file_path = slide_file_path.resolve().relative_to(Path.cwd().resolve())
                slide_file_path_str = str(rel_slide_file_path)
            except ValueError:
                slide_file_path_str = str(slide_file_path.resolve())
            
            # Record slide data for later DB insert after optional PDF extraction
            slide_data_for_pdf.append(
                {
                    "slide_id": slide_id,
                    "deck_id": deck_id,
                    "slide_index": slide_idx,
                    "title_header": title_header,
                    "plain_text": plain_text,
                    "summary": summary,
                    "thumbnail_path": thumbnail_path_str,
                    "original_slide_position": slide_idx,
                    "slide_file_path": slide_file_path_str,
                    "complexity_score": complexity_score,
                    "requires_pdf": requires_pdf,
                    "vector_id": vector_id,
                }
            )
            
            # Collect data for LightRAG batch insert
            if settings.lightrag_enabled:
                lightrag_documents.append({
                    'text': embedding_input,
                    'id': slide_id,
                    'metadata': {
                        'deck_id': deck_id,
                        'deck_filename': file_path.name,
                        'slide_index': slide_idx,
                        'title': title_header or f"Slide {slide_idx + 1}",
                    }
                })
            
            logger.debug(f"Slide {slide_idx + 1} processed successfully")
        
        # Convert entire deck to PDF and extract per-slide PDFs
        pdf_path = deck_pdf_path
        if settings.pdf_conversion_enabled and slide_data_for_pdf and pdf_path:
            try:
                    for data in slide_data_for_pdf:
                        slide_id = data["slide_id"]
                        slide_idx = data["slide_index"]

                        pdf_output_path = settings.slides_pdf_dir / f"{slide_id}.pdf"
                        success = pdf_processor.extract_pdf_page(
                            pdf_path, slide_idx, pdf_output_path
                        )

                        if success:
                            data["slide_pdf_path"] = str(pdf_output_path)
            except Exception as e:
                logger.error(f"PDF conversion failed: {e}")

        # Insert slide rows now that we have optional PDF paths
        for data in slide_data_for_pdf:
            db.insert_slide_with_id(
                slide_id=data["slide_id"],
                deck_id=data["deck_id"],
                slide_index=data["slide_index"],
                title_header=data["title_header"],
                plain_text=data["plain_text"],
                summary=data["summary"],
                thumbnail_path=data["thumbnail_path"],
                original_slide_position=data["original_slide_position"],
                slide_file_path=data["slide_file_path"],
                slide_pdf_path=data.get("slide_pdf_path"),
                requires_pdf=data["requires_pdf"],
                complexity_score=data["complexity_score"],
            )

        # Insert all slides into LightRAG in batch
        if settings.lightrag_enabled and lightrag_documents:
            try:
                logger.info(f"Inserting {len(lightrag_documents)} slides into LightRAG")
                lightrag_client.insert_documents_batch(lightrag_documents)
                logger.info("LightRAG batch insert complete")
            except Exception as e:
                logger.error(f"Error inserting documents into LightRAG: {e}")
                # Continue anyway as metadata is in PostgreSQL
        
        logger.info(f"Ingestion complete: {file_path} (deck_id: {deck_id})")
        return deck_id
    
    @staticmethod
    def ingest_folder(
        folder_path: Path,
        recursive: bool = True,
        uploader: Optional[str] = None,
        session_id: Optional[str] = None
    ) -> List[str]:
        """
        Ingest all PowerPoint files in a folder.
        
        Args:
            folder_path: Path to the folder
            recursive: Whether to search recursively
            uploader: Optional uploader name
            session_id: Optional session ID for audit logging
            
        Returns:
            List of deck_ids for successfully ingested files
        """
        folder_path = Path(folder_path).resolve()
        
        if not folder_path.exists():
            logger.error(f"Folder not found: {folder_path}")
            raise FileNotFoundError(f"Folder not found: {folder_path}")
        
        if not folder_path.is_dir():
            logger.error(f"Not a directory: {folder_path}")
            raise ValueError(f"Not a directory: {folder_path}")
        
        logger.info(f"Searching for .pptx files in: {folder_path} (recursive={recursive})")
        
        # Find all .pptx files
        if recursive:
            pptx_files = list(folder_path.rglob("*.pptx"))
        else:
            pptx_files = list(folder_path.glob("*.pptx"))
        
        # Filter out temporary files (starting with ~$)
        pptx_files = [f for f in pptx_files if not f.name.startswith("~$")]
        
        logger.info(f"Found {len(pptx_files)} .pptx files")
        
        deck_ids = []
        for pptx_file in pptx_files:
            try:
                deck_id = IngestEngine.ingest_file(pptx_file, uploader, session_id)
                if deck_id:
                    deck_ids.append(deck_id)
            except Exception as e:
                logger.error(f"Error ingesting {pptx_file}: {e}")
                # Continue with other files
                continue
        
        logger.info(f"Folder ingestion complete: {len(deck_ids)} new decks ingested")
        return deck_ids


# Global ingest engine instance
ingest_engine = IngestEngine()
